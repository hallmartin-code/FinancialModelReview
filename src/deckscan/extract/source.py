"""Turn an input file into content blocks Claude can read.

PDFs are sent as native document blocks so charts, tables, and image-only slides
are read from the page itself rather than from a lossy text dump. PPTX and
spreadsheets have no native block type, so they are flattened to text with
slide/cell markers that survive into the locators.

Nothing here interprets meaning, and nothing here raises on bad input: an
unreadable source degrades to an empty block list plus a methodology note, which
the pipeline turns into a recorded gap.
"""

from __future__ import annotations

import base64
import csv
import io
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from deckscan.config import Settings

#: Anthropic caps a request at 32MB; stay under it after base64 expansion (~1.34x).
MAX_NATIVE_PDF_BYTES = 20 * 1024 * 1024
#: Native PDF input is capped at 100 pages for 200k-context models; stay well inside.
MAX_NATIVE_PDF_PAGES = 100
#: Hard ceiling on inlined text, so a pathological file cannot blow up the request.
MAX_TEXT_CHARS = 400_000


@dataclass
class SourceContent:
    """One prepared input: what to send, and what to say about how it was read."""

    label: str
    """Human phrase used in the prompt, e.g. 'pitch deck (PDF)'."""
    blocks: list[dict[str, Any]] = field(default_factory=list)
    inline_text: str | None = None
    methodology: list[str] = field(default_factory=list)
    failed: bool = False
    """True when nothing readable came out of the file."""

    @property
    def is_empty(self) -> bool:
        return not self.blocks and not (self.inline_text or "").strip()


def _truncate(text: str) -> tuple[str, bool]:
    if len(text) <= MAX_TEXT_CHARS:
        return text, False
    return text[:MAX_TEXT_CHARS] + "\n\n[source truncated]", True


def prepare_deck(path: Path, settings: Settings) -> SourceContent:
    """Prepare a .pdf or .pptx deck."""
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _prepare_pdf(path, settings)
    if suffix == ".pptx":
        return _prepare_pptx(path)
    if suffix == ".docx":
        return _prepare_docx(path)
    return SourceContent(
        label="pitch deck",
        methodology=[f"Unsupported deck type {suffix!r}; nothing was read."],
        failed=True,
    )


def prepare_model(path: Path) -> SourceContent:
    """Prepare a .xlsx, .xlsm or .csv financial model."""
    suffix = path.suffix.lower()
    if suffix in {".xlsx", ".xlsm"}:
        return _prepare_xlsx(path)
    if suffix == ".csv":
        return _prepare_csv(path)
    return SourceContent(
        label="financial model",
        methodology=[f"Unsupported model type {suffix!r}; nothing was read."],
        failed=True,
    )


# --- PDF ---------------------------------------------------------------------


def _pdf_text(path: Path) -> tuple[str, int, int]:
    """Extracted text with page markers, page count, and pages that had no text."""
    import pdfplumber

    chunks: list[str] = []
    pages = 0
    empty = 0
    with pdfplumber.open(str(path)) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            pages += 1
            text = (page.extract_text() or "").strip()
            if not text:
                empty += 1
                continue
            chunks.append(f"[p.{index}]\n{text}")
    return "\n\n".join(chunks), pages, empty


def _prepare_pdf(path: Path, settings: Settings) -> SourceContent:
    content = SourceContent(label="pitch deck (PDF)")

    try:
        text, pages, empty_pages = _pdf_text(path)
    except Exception as exc:  # pdfplumber raises a wide range on damaged files
        text, pages, empty_pages = "", 0, 0
        content.methodology.append(f"PDF text layer unreadable ({type(exc).__name__}).")

    size = path.stat().st_size
    native = size <= MAX_NATIVE_PDF_BYTES and 0 < pages <= MAX_NATIVE_PDF_PAGES

    if native:
        data = base64.standard_b64encode(path.read_bytes()).decode("ascii")
        content.blocks.append(
            {
                "type": "document",
                "source": {"type": "base64", "media_type": "application/pdf", "data": data},
                "title": path.name,
            }
        )
        content.methodology.append(
            f"Deck read as a native PDF ({pages} pages), so charts and image-only "
            "pages were read from the page itself."
        )
        if empty_pages:
            content.methodology.append(
                f"{empty_pages} of {pages} pages carry no text layer and were read visually."
            )
        return content

    if pages > MAX_NATIVE_PDF_PAGES:
        content.methodology.append(
            f"Deck has {pages} pages, above the {MAX_NATIVE_PDF_PAGES}-page native limit; "
            "extracted text was sent instead and any image-only pages were not read."
        )
    elif size > MAX_NATIVE_PDF_BYTES:
        content.methodology.append(
            f"Deck is {size // (1024 * 1024)}MB, above the "
            f"{MAX_NATIVE_PDF_BYTES // (1024 * 1024)}MB native limit; extracted text "
            "was sent instead and any image-only pages were not read."
        )

    if not text.strip():
        content.failed = True
        content.methodology.append("No text could be extracted from the deck.")
        return content

    body, truncated = _truncate(text)
    content.inline_text = body
    if truncated:
        content.methodology.append("Deck text was truncated to fit the request.")
    if empty_pages and pages:
        content.methodology.append(
            f"{empty_pages} of {pages} pages carry no text layer and were not read. "
            f"(Threshold: {settings.extraction.min_chars_per_page} chars per page.)"
        )
    return content


# --- PPTX --------------------------------------------------------------------


def _shape_text(shape: Any) -> list[str]:
    """Text from one shape, recursing into groups and reading table cells."""
    out: list[str] = []
    if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
        for child in shape.shapes:
            out.extend(_shape_text(child))
        return out
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                out.append(line)
        return out
    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            line = "".join(run.text for run in para.runs).strip()
            if line:
                out.append(line)
        return out
    text = getattr(shape, "text", "")
    if isinstance(text, str) and text.strip():
        out.append(text.strip())
    return out


def _prepare_pptx(path: Path) -> SourceContent:
    content = SourceContent(label="pitch deck (PowerPoint)")
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
    except Exception as exc:
        content.failed = True
        content.methodology.append(f"PowerPoint file could not be opened ({type(exc).__name__}).")
        return content

    chunks: list[str] = []
    picture_only = 0
    for index, slide in enumerate(prs.slides, start=1):
        lines: list[str] = []
        for shape in slide.shapes:
            lines.extend(_shape_text(shape))
        notes = ""
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        if not lines and not notes:
            picture_only += 1
            continue
        body = "\n".join(lines)
        if notes:
            body += f"\n[speaker notes] {notes}"
        chunks.append(f"[slide {index}]\n{body}")

    if not chunks:
        content.failed = True
        content.methodology.append("No text could be extracted from any slide.")
        return content

    text, truncated = _truncate("\n\n".join(chunks))
    content.inline_text = text
    content.methodology.append(f"Deck read as PowerPoint text ({len(chunks)} slides with text).")
    if picture_only:
        content.methodology.append(
            f"{picture_only} slides contain only images and could not be read; "
            "supply the deck as a PDF to have those pages read visually."
        )
    if truncated:
        content.methodology.append("Deck text was truncated to fit the request.")
    return content


# --- DOCX --------------------------------------------------------------------


def _prepare_docx(path: Path) -> SourceContent:
    """Walk the document body in order, splitting on explicit page breaks.

    Decks are sometimes written as Word documents. Page breaks are the closest
    thing the format has to slide boundaries, so they become the locator unit.
    """
    content = SourceContent(label="pitch deck (Word)")
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        document = Document(str(path))
    except Exception as exc:
        content.failed = True
        content.methodology.append(f"Word document could not be opened ({type(exc).__name__}).")
        return content

    def has_page_break(element: Any) -> bool:
        if any(node.get(qn("w:type")) == "page" for node in element.iter(qn("w:br"))):
            return True
        # Word writes this where it laid out a page.
        return next(element.iter(qn("w:lastRenderedPageBreak")), None) is not None

    pages: list[list[str]] = [[]]
    for child in document.element.body.iterchildren():
        tag = child.tag.split("}")[-1]
        if tag == "p":
            text = Paragraph(child, document).text.strip()
            if text:
                pages[-1].append(text)
            if has_page_break(child):
                pages.append([])
        elif tag == "tbl":
            for row in Table(child, document).rows:
                cells = [cell.text.strip() for cell in row.cells]
                line = " | ".join(c for c in cells if c)
                if line:
                    pages[-1].append(line)

    chunks = [
        f"[p.{index}]\n" + "\n".join(lines)
        for index, lines in enumerate((p for p in pages if p), start=1)
    ]
    if not chunks:
        content.failed = True
        content.methodology.append("No text could be extracted from the document.")
        return content

    text, truncated = _truncate("\n\n".join(chunks))
    content.inline_text = text
    content.methodology.append(
        f"Deck read as Word text ({len(chunks)} page-equivalents); any embedded "
        "images were not read."
    )
    if truncated:
        content.methodology.append("Deck text was truncated to fit the request.")
    return content


# --- spreadsheets ------------------------------------------------------------


def _column_letter(index: int) -> str:
    """1 -> A, 27 -> AA."""
    letters = ""
    while index > 0:
        index, remainder = divmod(index - 1, 26)
        letters = chr(65 + remainder) + letters
    return letters


@dataclass
class _SheetRead:
    """One tab as it came off the workbook, before any budgeting."""

    title: str
    hidden: bool
    rows: list[str]
    numeric_cells: int

    def block(self, rows: list[str] | None = None) -> str:
        body = self.rows if rows is None else rows
        return f"[sheet {self.title}]\n" + "\n".join(body)

    @property
    def size(self) -> int:
        return len(self.block())


def _read_sheet(sheet: Any) -> tuple[list[str], int]:
    """One tab's populated rows as ``A1=value`` lines, and how many held numbers."""
    rows: list[str] = []
    numeric = 0
    for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
        cells: list[str] = []
        for col, value in enumerate(row, start=1):
            if value is None or str(value).strip() == "":
                continue
            if isinstance(value, int | float) and not isinstance(value, bool):
                numeric += 1
            cells.append(f"{_column_letter(col)}{row_index}={value}")
        if cells:
            rows.append("  ".join(cells))
    return rows, numeric


def _fair_shares(sizes: list[int], budget: int) -> list[int]:
    """Split ``budget`` across ``sizes``, max-min fair.

    A tab that fits inside an equal share keeps all of its text, and what it does
    not use is redistributed to the larger tabs. This is what stops a long first
    sheet from spending the whole request and pushing later tabs out entirely:
    every tab in the workbook reaches Claude, even if a huge one arrives clipped.
    """
    shares = [0] * len(sizes)
    pending = set(range(len(sizes)))
    left = budget
    while pending:
        share = left // len(pending)
        fitting = {index for index in pending if sizes[index] <= share}
        if not fitting:
            for index in pending:
                shares[index] = share
            break
        for index in fitting:
            shares[index] = sizes[index]
            left -= sizes[index]
        pending -= fitting
    return shares


#: Room held back from a clipped tab's allowance for the "further rows" marker.
_MARKER_RESERVE = 80


def _trim_to(sheet: _SheetRead, allowance: int) -> tuple[str, int]:
    """The tab's block clipped to ``allowance`` characters on whole-row boundaries."""
    kept: list[str] = []
    used = len(f"[sheet {sheet.title}]")
    for row in sheet.rows:
        if used + 1 + len(row) > allowance - _MARKER_RESERVE:
            break
        kept.append(row)
        used += 1 + len(row)
    if not kept:  # never emit a bare header; one row is what proves the tab exists
        kept = sheet.rows[:1]

    def render(shown: list[str]) -> str:
        dropped = len(sheet.rows) - len(shown)
        marker = [f"[... {dropped} further rows on this tab not shown]"] if dropped else []
        return sheet.block([*shown, *marker])

    # A long first row or a long marker can still overshoot; give rows back until it fits.
    while len(kept) > 1 and len(render(kept)) > allowance:
        kept.pop()
    return render(kept), len(kept)


def _formula_cell_counts(path: Path, titles: set[str]) -> dict[str, int]:
    """Formula cells on the named tabs, read as formulas rather than cached values.

    Only called for tabs that produced no numbers at all, to tell "this tab is
    prose" apart from "this tab is formulas whose results were never cached".
    """
    if not titles:
        return {}
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), data_only=False, read_only=True)
    except Exception:
        return {}
    counts: dict[str, int] = {}
    try:
        for sheet in workbook.worksheets:
            if sheet.title not in titles:
                continue
            found = sum(
                1
                for row in sheet.iter_rows(values_only=True)
                for value in row
                if isinstance(value, str) and value.startswith("=")
            )
            if found:
                counts[sheet.title] = found
    except Exception:
        return counts
    finally:
        workbook.close()
    return counts


def _prepare_xlsx(path: Path) -> SourceContent:
    """Read every tab in the workbook, budgeting the request across all of them."""
    content = SourceContent(label="financial model (Excel)")
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), data_only=True, read_only=True)
    except Exception as exc:
        content.failed = True
        content.methodology.append(f"Workbook could not be opened ({type(exc).__name__}).")
        return content

    sheets: list[_SheetRead] = []
    try:
        for sheet in workbook.worksheets:
            rows, numeric = _read_sheet(sheet)
            sheets.append(
                _SheetRead(
                    title=sheet.title,
                    hidden=getattr(sheet, "sheet_state", "visible") != "visible",
                    rows=rows,
                    numeric_cells=numeric,
                )
            )
    finally:
        workbook.close()

    populated = [tab for tab in sheets if tab.rows]
    if not populated:
        content.failed = True
        content.methodology.append(
            f"The workbook has {len(sheets)} tab(s) and no populated cells in any of them."
        )
        return content

    # Blocks are joined by a blank line, so the separators come out of the budget too.
    budget = MAX_TEXT_CHARS - 2 * (len(populated) - 1)
    shares = _fair_shares([tab.size for tab in populated], budget)
    blocks: list[str] = []
    clipped: list[str] = []
    for tab, allowance in zip(populated, shares, strict=True):
        if tab.size <= allowance:
            blocks.append(tab.block())
            continue
        block, shown = _trim_to(tab, allowance)
        blocks.append(block)
        clipped.append(f"{tab.title} ({shown:,} of {len(tab.rows):,} rows)")

    content.inline_text = "\n\n".join(blocks)

    names = ", ".join(tab.title for tab in populated)
    content.methodology.append(
        f"Model read from all {len(populated)} populated tab(s) of {len(sheets)}: {names}. "
        "Every value carries its sheet and cell reference."
    )
    empty = [tab.title for tab in sheets if not tab.rows]
    if empty:
        content.methodology.append(f"Tabs with no populated cells: {', '.join(empty)}.")
    hidden = [tab.title for tab in populated if tab.hidden]
    if hidden:
        content.methodology.append(f"Hidden tabs were read as well: {', '.join(hidden)}.")
    if clipped:
        content.methodology.append(
            "These tabs were too long to send whole and were clipped to fit: "
            + "; ".join(clipped)
            + "."
        )

    # A tab of formulas whose results Excel never cached reads as labels with no
    # figures. Say so, rather than letting it look like a tab that states nothing.
    uncached = _formula_cell_counts(
        path, {tab.title for tab in populated if tab.numeric_cells == 0}
    )
    if uncached:
        detail = ", ".join(f"{title} ({count:,} cells)" for title, count in uncached.items())
        content.methodology.append(
            f"These tabs hold formulas with no cached results, so their figures could not "
            f"be read: {detail}. Open the workbook in Excel and re-save it to include them."
        )
    return content


def _prepare_csv(path: Path) -> SourceContent:
    content = SourceContent(label="financial model (CSV)")
    try:
        raw = path.read_text(encoding="utf-8-sig", errors="replace")
        rows = list(csv.reader(io.StringIO(raw)))
    except Exception as exc:
        content.failed = True
        content.methodology.append(f"CSV could not be read ({type(exc).__name__}).")
        return content

    lines: list[str] = []
    for row_index, row in enumerate(rows, start=1):
        cells = [
            f"{_column_letter(col)}{row_index}={value.strip()}"
            for col, value in enumerate(row, start=1)
            if value.strip()
        ]
        if cells:
            lines.append("  ".join(cells))

    if not lines:
        content.failed = True
        content.methodology.append("The CSV contained no populated cells.")
        return content

    text, truncated = _truncate(f"[sheet {path.stem}]\n" + "\n".join(lines))
    content.inline_text = text
    content.methodology.append("Model read from CSV; every value carries its cell reference.")
    if truncated:
        content.methodology.append("Model text was truncated to fit the request.")
    return content
